from pathlib import Path
from shutil import copy2
from pptx import Presentation

ppt_path = Path('/Users/macstudio/Documents/Intern_Peking University_supu/2026年04月实习生汇报/2026年04月实习生汇报-江浩华_version2.pptx')
backup_path = ppt_path.with_name(ppt_path.stem + '_before_notes_backup.pptx')

notes = [
    "大家好，我是江浩华，下面汇报我在2026年4月围绕CPU并行整体刚度矩阵组装所开展的工作。本月汇报的核心，不是单独展示某一个算法实现，而是汇报这一方向已经从前期原型探索推进到可复现实验平台建设与阶段验证。",
    "这一部分我主要从三个层面展开。第一，为什么项目需要从前期的算法调研转向平台化建设；第二，本月在统一主线、实验流程和结果产物方面完成了哪些关键工作；第三，在真实工程网格上，我们已经得到哪些有代表性的阶段性结论。",
    "先介绍项目背景。1月阶段我们已经完成了多类并行组装策略调研，但当时整体叙事仍然偏向GPU-first。进入4月以后，项目目标发生了明显变化，不再只是继续做CPU原型，而是转向CPU可复现实验平台建设。对应到本月目标，主要有四项：第一，固定唯一CPU主线目录并统一核心文档；第二，跑通真实工程网格3d-WindTurbineHub的实验矩阵；第三，形成CSV、JSON、Markdown、PNG、SVG五类结构化结果；第四，基于真实数据形成阶段性算法判断。",
    "这一页汇报本月在平台建设方面完成的主要内容。第一是主线收敛，也就是把CPU并行唯一主线目录固定为cpu_parallel_stiffness_assembly，同时把历史GPU资产隔离为legacy参考，并统一README、requirements和handoff文档口径。第二是实验框架统一，包括统一算法工厂、benchmark命令行入口、规则网格与inp解析、CSR和scatter plan，使不同算法可以在同一口径下公平比较。整体上看，本月工作的重点不是新增零散功能，而是完成主线统一、实验补齐、结果固化和结论提炼这几项平台化建设。",
    "这一页主要讲实验设计和验证流程。本月实验强调统一流程，也就是统一输入、统一算法和统一输出。在输入上，一方面保留规则网格用于回归测试，另一方面重点引入真实工程网格3d-WindTurbineHub。在规模上，这个工程算例已经达到22万级节点、111万级单元、68万级自由度以及2750万级NNZ，因此已经能够支撑工程级比较，而不是停留在小规模样例层面。也正因为有了这样的统一流程和统一算例，后面的结果才具有可复现和可比较的意义。",
    "下面开始汇报真实工程网格上的结果。先看simplified kernel场景。在当前设置下，最优算法是row_owner，它在14线程下达到了55.141毫秒，对串行基线实现了3.695倍加速。与此同时，private_csr达到71.211毫秒、2.861倍加速；atomic在额外内存几乎为零的前提下达到80.560毫秒。这个结果说明，在simplified场景下，owner-computes，也就是row_owner这一路线，在高线程条件下优势最明显。",
    "接着看更接近真实计算代价的physics_tet4场景。在这个kernel下，第一梯队已经比较清楚地收敛为private_csr和row_owner。其中最优结果是private_csr在8线程达到119.566毫秒、4.686倍加速。row_owner在8线程的结果几乎持平，为120.444毫秒、4.652倍；atomic也达到123.174毫秒、4.549倍。也就是说，在更复杂的kernel下，线程私有缓冲归并路线体现出更好的稳定性，而row_owner依然保持非常强的竞争力。",
    "这一页对不同算法做一个阶段性归纳。当前最值得继续深入的两类路线已经比较清楚，就是row_owner和private_csr。其中，row_owner在simplified场景下表现最强，而private_csr在physics_tet4场景下略优。atomic虽然不是绝对最快，但由于额外内存几乎为零，所以工程性价比非常高。coloring可以稳定运行，但已经不是当前领先路线；coo_sort_reduce虽然正确性成立，但速度和资源代价都不具备优势，更适合作为研究对照组，而不适合作为后续主推方向。",
    "以上就是我本月的汇报内容。总体来看，本月工作的意义不在于单独实现了某一个算法，而在于把CPU并行整体刚度矩阵组装研究推进成了一套可验证、可比较、可复现的工程级实验平台。后续我会继续围绕真实工程网格上的线程扫描、跨平台复现，以及row_owner、private_csr和atomic这几条重点路线做进一步分析。感谢老师和各位的指正。",
]

if not backup_path.exists():
    copy2(ppt_path, backup_path)

prs = Presentation(str(ppt_path))
assert len(prs.slides) == len(notes), f'slides={len(prs.slides)} notes={len(notes)}'

for slide, note_text in zip(prs.slides, notes):
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = note_text

prs.save(str(ppt_path))
print(str(ppt_path))
print(str(backup_path))
