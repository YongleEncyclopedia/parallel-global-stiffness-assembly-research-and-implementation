#!/usr/bin/env python3
"""V2: readable comparison charts for five CPU parallel algorithms.
Generates detailed + concise figures and a one-slide summary with three simplified charts.
"""
from pathlib import Path
import csv, subprocess, sys
import numpy as np
import matplotlib.pyplot as plt
from matplotlib import font_manager
from matplotlib.patches import FancyBboxPatch, Rectangle

ROOT = Path(__file__).resolve().parents[1]
CSV = ROOT / "results/2026-04-22/csv/windhub_simplified.csv"
OUT = ROOT / "presentation_assets/2026-04-cpu-algorithm-comparison-v2"
OUT.mkdir(parents=True, exist_ok=True)
for fp in ["/System/Library/Fonts/PingFang.ttc", "/System/Library/Fonts/STHeiti Medium.ttc", "/System/Library/Fonts/Supplemental/Songti.ttc"]:
    if Path(fp).exists():
        font_manager.fontManager.addfont(fp)
        plt.rcParams["font.family"] = font_manager.FontProperties(fname=fp).get_name(); break
plt.rcParams.update({"axes.unicode_minus":False,"font.size":15,"axes.titlesize":20,"axes.labelsize":16,"xtick.labelsize":14,"ytick.labelsize":14,"legend.fontsize":13})

ALGOS=[("cpu_atomic","Atomic","#2563EB"),("cpu_private_csr","Private CSR","#16A34A"),("cpu_coo_sort_reduce","COO\nSort-Reduce","#F97316"),("cpu_graph_coloring","Graph\nColoring","#8B5CF6"),("cpu_row_owner","Row Owner","#E11D48")]
LABEL={a:b for a,b,c in ALGOS}; COLOR={a:c for a,b,c in ALGOS}; ORDER=[a for a,b,c in ALGOS]
BG="#F8FAFC"; INK="#0F172A"; MUTED="#475569"; GRID="#E2E8F0"

def read_rows():
    rows=[]
    with open(CSV,encoding='utf-8') as f:
        for r in csv.DictReader(f):
            if r['algorithm'] in ORDER:
                d=dict(r)
                for k in ['threads','assembly_mean_ms','total_mean_ms','speedup','efficiency','rel_l2','max_abs','extra_memory_bytes','peak_rss_mb']:
                    d[k]=float(d[k])
                d['threads']=int(d['threads']); d['extra_memory_gb']=d['extra_memory_bytes']/1024**3; d['peak_rss_gb']=d['peak_rss_mb']/1024
                rows.append(d)
    return rows
ROWS=read_rows(); THREADS=sorted({r['threads'] for r in ROWS})
def rows_alg(a): return sorted([r for r in ROWS if r['algorithm']==a], key=lambda r:r['threads'])
def best(a): return max(rows_alg(a), key=lambda r:r['speedup'])

def fig_title(fig,title,subtitle,dark=False):
    fig.text(.055,.94,title,fontsize=30,fontweight='bold',color=('white' if dark else INK))
    fig.text(.055,.898,subtitle,fontsize=15,color=('#CBD5E1' if dark else MUTED))
def style(ax,title=None):
    ax.set_facecolor('white'); ax.grid(True,color=GRID,lw=1); ax.set_axisbelow(True)
    if title: ax.set_title(title,fontweight='bold',pad=14)
    for s in ax.spines.values(): s.set_color('#CBD5E1')
    ax.tick_params(colors='#334155')
def save(fig,name):
    png=OUT/f'{name}.png'; svg=OUT/f'{name}.svg'
    fig.savefig(png,dpi=220,facecolor=fig.get_facecolor()); fig.savefig(svg,facecolor=fig.get_facecolor()); plt.close(fig); return [png,svg]

def add_common_legend(fig, handles, labels):
    fig.legend(handles, labels, loc='upper center', bbox_to_anchor=(0.5,0.84), ncol=5, frameon=True, facecolor='white', edgecolor='#CBD5E1', fontsize=13.5)

def detailed_correctness():
    fig=plt.figure(figsize=(16,9),facecolor=BG); fig_title(fig,'Correctness Comparison / 正确性对比','Actual WindTurbineHub simplified benchmark; errors are measured against serial CSR assembly')
    ax1=fig.add_axes([.075,.18,.40,.56]); ax2=fig.add_axes([.555,.18,.38,.56])
    handles=[]; labels=[]
    for ax,title,field,ylabel,band in [(ax1,'Relative L2 matrix-norm error','rel_l2','rel_L2 error',(1e-18,1e-12)),(ax2,'Max absolute matrix-entry error','max_abs','max_abs error',(1e-18,1e-10))]:
        style(ax,title); ax.set_yscale('log'); ax.axhspan(band[0],band[1],color='#DCFCE7',alpha=.55,zorder=0)
        for alg,lab,c in ALGOS:
            rs=rows_alg(alg); line=ax.plot([r['threads'] for r in rs],[max(r[field],1e-18) for r in rs],marker='o',ms=6.5,lw=2.8,color=c,label=lab.replace('\n',' '))[0]
            if ax is ax1: handles.append(line); labels.append(lab.replace('\n',' '))
        ax.set_xlabel('Threads'); ax.set_ylabel(ylabel); ax.set_xticks([1,2,4,6,8,10,12,14])
    add_common_legend(fig,handles,labels)
    fig.text(.075,.085,'Conclusion: all methods remain at round-off level; zero-error cases are shown at 1e-18 for log-scale visibility.',fontsize=16,color=MUTED)
    return save(fig,'01_detailed_correctness_v2')

def detailed_efficiency():
    fig=plt.figure(figsize=(16,9),facecolor=BG); fig_title(fig,'Efficiency Comparison / 效率对比','Assembly speedup and parallel efficiency over 1–14 CPU threads')
    ax1=fig.add_axes([.075,.18,.40,.56]); ax2=fig.add_axes([.555,.18,.38,.56]); handles=[]; labels=[]
    style(ax1,'Assembly speedup vs serial baseline'); style(ax2,'Parallel efficiency = speedup / threads')
    ideal=ax1.plot(THREADS,THREADS,'--',lw=2,color='#94A3B8',label='Ideal linear')[0]
    handles.append(ideal); labels.append('Ideal linear')
    for alg,lab,c in ALGOS:
        rs=rows_alg(alg); x=[r['threads'] for r in rs]
        h=ax1.plot(x,[r['speedup'] for r in rs],marker='o',ms=6.5,lw=2.8,color=c,label=lab.replace('\n',' '))[0]
        ax2.plot(x,[r['efficiency'] for r in rs],marker='o',ms=6.5,lw=2.8,color=c)
        handles.append(h); labels.append(lab.replace('\n',' '))
    ax1.set_xlabel('Threads'); ax1.set_ylabel('Speedup ×'); ax1.set_ylim(0,14.5); ax1.set_xticks([1,2,4,6,8,10,12,14])
    ax2.set_xlabel('Threads'); ax2.set_ylabel('Efficiency'); ax2.set_ylim(0,1.05); ax2.set_xticks([1,2,4,6,8,10,12,14])
    fig.legend(handles,labels,loc='upper center',bbox_to_anchor=(0.5,0.84),ncol=6,frameon=True,facecolor='white',edgecolor='#CBD5E1',fontsize=12.8)
    fig.text(.075,.085,'Conclusion: Row Owner is fastest; Private CSR is competitive; COO Sort-Reduce is a slow research baseline.',fontsize=16,color=MUTED)
    return save(fig,'02_detailed_efficiency_v2')

def detailed_memory():
    fig=plt.figure(figsize=(16,9),facecolor=BG); fig_title(fig,'Memory Footprint Comparison / 内存占用对比','Extra algorithmic memory and observed peak RSS across thread counts')
    ax1=fig.add_axes([.075,.18,.40,.56]); ax2=fig.add_axes([.555,.18,.38,.56]); handles=[]; labels=[]
    style(ax1,'Extra memory requested by algorithm'); style(ax2,'Peak RSS during benchmark')
    for alg,lab,c in ALGOS:
        rs=rows_alg(alg); x=[r['threads'] for r in rs]
        h=ax1.plot(x,[r['extra_memory_gb'] for r in rs],marker='o',ms=6.5,lw=2.8,color=c,label=lab.replace('\n',' '))[0]
        ax2.plot(x,[r['peak_rss_gb'] for r in rs],marker='o',ms=6.5,lw=2.8,color=c)
        handles.append(h); labels.append(lab.replace('\n',' '))
    ax1.set_xlabel('Threads'); ax1.set_ylabel('Extra memory (GB)'); ax1.set_xticks([1,2,4,6,8,10,12,14])
    ax2.set_xlabel('Threads'); ax2.set_ylabel('Peak RSS (GB)'); ax2.set_xticks([1,2,4,6,8,10,12,14])
    add_common_legend(fig,handles,labels)
    fig.text(.075,.085,'Conclusion: Atomic/Graph Coloring are memory-light; Private CSR grows with threads; COO/Owner keep large buffers.',fontsize=16,color=MUTED)
    return save(fig,'03_detailed_memory_v2')

# concise standalone charts
def multiline_labels(): return [LABEL[a] for a in ORDER]
def simple_style(ax,title):
    style(ax,title); ax.tick_params(axis='x',labelsize=13); ax.tick_params(axis='y',labelsize=13)
def simple_correct(ax):
    vals=[max(max(r['rel_l2'],1e-18) for r in rows_alg(a)) for a in ORDER]
    cols=[COLOR[a] for a in ORDER]; simple_style(ax,'Correctness: worst relative L2 error')
    bars=ax.bar(range(5),vals,color=cols,width=.68); ax.set_yscale('log'); ax.set_ylim(1e-18,1e-14); ax.set_ylabel('Worst rel_L2'); ax.set_xticks(range(5),multiline_labels())
    ax.axhspan(1e-18,1e-12,color='#DCFCE7',alpha=.55,zorder=0)
    for b,v in zip(bars,vals): ax.text(b.get_x()+b.get_width()/2,v*1.55,f'{v:.1e}',ha='center',va='bottom',fontsize=14,color=INK,fontweight='bold')
def simple_eff(ax):
    vals=[]; th=[]
    for a in ORDER:
        b=best(a); vals.append(b['speedup']); th.append(b['threads'])
    simple_style(ax,'Efficiency: best observed speedup')
    bars=ax.bar(range(5),vals,color=[COLOR[a] for a in ORDER],width=.68); ax.set_ylabel('Speedup ×'); ax.set_ylim(0,max(vals)*1.28); ax.set_xticks(range(5),multiline_labels())
    for b,v,t in zip(bars,vals,th): ax.text(b.get_x()+b.get_width()/2,max(v+.10,.18),f'{v:.2f}×\n@{t}T',ha='center',va='bottom',fontsize=14,color=INK,fontweight='bold')
def simple_mem(ax):
    vals=[best(a)['extra_memory_gb'] for a in ORDER]
    simple_style(ax,'Memory: extra memory at best speedup')
    bars=ax.bar(range(5),vals,color=[COLOR[a] for a in ORDER],width=.68); ax.set_ylabel('Extra memory (GB)'); ax.set_ylim(0,max(vals)*1.25+.15); ax.set_xticks(range(5),multiline_labels())
    for b,v in zip(bars,vals): ax.text(b.get_x()+b.get_width()/2,v+.07,f'{v:.2f}',ha='center',va='bottom',fontsize=14,color=INK,fontweight='bold')

def standalone_simple():
    outs=[]
    for name,title,func in [('04_simple_correctness_v2','Correctness / 正确性简洁版',simple_correct),('05_simple_efficiency_v2','Efficiency / 效率简洁版',simple_eff),('06_simple_memory_v2','Memory / 内存简洁版',simple_mem)]:
        fig=plt.figure(figsize=(16,9),facecolor=BG); fig_title(fig,title,'Five CPU parallel algorithms · actual WindTurbineHub simplified benchmark')
        ax=fig.add_axes([.09,.18,.84,.62]); func(ax); outs += save(fig,name)
    return outs

# summary slide: redraw compact horizontal-bar charts directly, no tiny axes labels
def panel(ax,title,subtitle):
    ax.set_axis_off(); ax.add_patch(FancyBboxPatch((0,0),1,1,boxstyle='round,pad=0.018,rounding_size=0.04',facecolor='white',edgecolor='#93A4BC',lw=1.4))
    ax.text(.055,.93,title,fontsize=20,fontweight='bold',color=INK,va='top'); ax.text(.055,.855,subtitle,fontsize=12.5,color=MUTED,va='top')

def draw_hbars(ax, labels, vals, cols, fmt, x0=.08, x1=.93, y0=.18, ygap=.12, log=False):
    maxv=max(vals) if not log else max(np.log10(v) for v in vals)
    minlog=min(np.log10(v) for v in vals) if log else 0
    for i,(lab,v,c) in enumerate(zip(labels,vals,cols)):
        y=.74-i*ygap;        ax.text(x0,y,lab.replace('\n',' '),fontsize=12.5,color=INK,ha='left',va='center',fontweight='bold')
        bx=x0+.36; bw=x1-bx
        ax.add_patch(Rectangle((bx,y-.025),bw,.05,facecolor='#E2E8F0',edgecolor='none'))
        frac=(np.log10(v)-minlog)/(maxv-minlog) if log and maxv>minlog else (v/maxv if maxv else 0)
        ax.add_patch(Rectangle((bx,y-.025),max(frac*bw,.012),.05,facecolor=c,edgecolor='none'))
        ax.text(x1,y,fmt(v),fontsize=12.3,color=INK,ha='right',va='center')

def summary_slide():
    fig=plt.figure(figsize=(16,9),facecolor='#07111F')
    bg=fig.add_axes([0,0,1,1]); bg.set_axis_off(); bg.set_xlim(0,1); bg.set_ylim(0,1)
    X,Y=np.meshgrid(np.linspace(0,1,500),np.linspace(0,1,300)); grad=np.zeros((300,500,3)); base=np.array([7,17,31])/255; grad[:]=base
    grad += np.array([37,99,235])/255*np.clip(1-((X-.18)**2+(Y-.9)**2)**0.5/.85,0,1)[...,None]*.28
    grad += np.array([225,29,72])/255*np.clip(1-((X-.92)**2+(Y-.12)**2)**0.5/.7,0,1)[...,None]*.18
    bg.imshow(np.clip(grad,0,1),extent=[0,1,0,1],origin='lower',aspect='auto')
    fig.text(.055,.93,'Five CPU Parallel Algorithms: Correctness · Efficiency · Memory',fontsize=27,fontweight='bold',color='white')
    fig.text(.055,.892,'Concise benchmark summary · 3d-WindTurbineHub / simplified kernel',fontsize=14,color='#CBD5E1')
    labels=multiline_labels(); cols=[COLOR[a] for a in ORDER]
    ax1=fig.add_axes([.055,.20,.285,.60]); panel(ax1,'Correctness','worst relative L2 error'); draw_hbars(ax1,labels,[max(max(r['rel_l2'],1e-18) for r in rows_alg(a)) for a in ORDER],cols,lambda v:f'{v:.1e}',log=True)
    ax2=fig.add_axes([.36,.20,.285,.60]); panel(ax2,'Efficiency','best observed speedup'); draw_hbars(ax2,labels,[best(a)['speedup'] for a in ORDER],cols,lambda v:f'{v:.2f}×')
    ax3=fig.add_axes([.665,.20,.285,.60]); panel(ax3,'Memory','extra GB at best speedup'); draw_hbars(ax3,labels,[best(a)['extra_memory_gb'] for a in ORDER],cols,lambda v:f'{v:.2f} GB')
    ax=fig.add_axes([.055,.065,.895,.085]); ax.set_axis_off(); ax.add_patch(FancyBboxPatch((0,0),1,1,boxstyle='round,pad=0.014,rounding_size=0.035',facecolor='#0F172A',edgecolor='#334155',lw=1.2,alpha=.94))
    ax.text(.025,.63,'Takeaway',fontsize=16.5,color='#93C5FD',fontweight='bold',va='center')
    ax.text(.15,.67,'Correctness: all methods pass.  Speed: Row Owner leads.',fontsize=15.2,color='white',va='center')
    ax.text(.15,.32,'Memory: Atomic / Graph Coloring are light; Private CSR / COO / Owner use extra buffers.',fontsize=15.2,color='white',va='center')
    return save(fig,'07_summary_slide_three_simple_charts_v2')

def summary_md():
    p=OUT/'data_summary_v2.md'; lines=[f'Data source: {CSV}\n\n','| Algorithm | Best speedup | Threads | Worst rel_L2 | Extra memory at best (GB) | Peak RSS at best (GB) |\n','|---|---:|---:|---:|---:|---:|\n']
    for a,lab,c in ALGOS:
        b=best(a); lines.append(f"| {lab.replace(chr(10),' ')} | {b['speedup']:.3f} | {b['threads']} | {max(r['rel_l2'] for r in rows_alg(a)):.3e} | {b['extra_memory_gb']:.3f} | {b['peak_rss_gb']:.2f} |\n")
    p.write_text(''.join(lines),encoding='utf-8'); return p

def pptx(summary_png):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        subprocess.check_call([sys.executable,'-m','pip','install','python-pptx','-q'])
        from pptx import Presentation
        from pptx.util import Inches
    prs=Presentation(); prs.slide_width=Inches(16); prs.slide_height=Inches(9)
    sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.shapes.add_picture(str(summary_png),0,0,width=prs.slide_width,height=prs.slide_height)
    out=OUT/'cpu_algorithm_comparison_summary_slide_v2.pptx'; prs.save(out); return out

if __name__=='__main__':
    outs=[]; outs+=detailed_correctness(); outs+=detailed_efficiency(); outs+=detailed_memory(); outs+=standalone_simple(); summary_png,summary_svg=summary_slide(); outs += [summary_png,summary_svg,summary_md(),pptx(summary_png)]
    print('Generated V2 comparison assets:')
    for p in outs: print(p)
