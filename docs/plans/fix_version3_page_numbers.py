from pptx import Presentation
from pathlib import Path
from pptx.enum.shapes import MSO_SHAPE_TYPE

ppt_path = Path('/Users/macstudio/Documents/Intern_Peking University_supu/2026年04月实习生汇报/2026年04月实习生汇报-江浩华_version3.pptx')
ppt = Presentation(str(ppt_path))
for idx, slide in enumerate(ppt.slides, start=1):
    # skip cover slide because it has no page footer
    if idx == 1:
        continue
    target = str(idx)
    footer_candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text'):
            continue
        txt = shape.text.strip()
        if not txt:
            continue
        if txt.isdigit() and shape.left > 8000000 and shape.top > 6200000:
            footer_candidates.append(shape)
    if footer_candidates:
        # right-most footer candidate
        footer = sorted(footer_candidates, key=lambda s: s.left)[-1]
        footer.text = target

ppt.save(str(ppt_path))
print(ppt_path)
