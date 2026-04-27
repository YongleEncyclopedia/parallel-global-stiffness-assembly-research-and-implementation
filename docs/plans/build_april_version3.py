from __future__ import annotations

import csv
import subprocess
import zipfile
from pathlib import Path
from shutil import copy2

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path('/Users/macstudio/Documents/Intern_Peking University_supu')
PPT_DIR = BASE_DIR / '2026年04月实习生汇报'
VERSION2 = PPT_DIR / '2026年04月实习生汇报-江浩华_version2.pptx'
VERSION3 = PPT_DIR / '2026年04月实习生汇报-江浩华_version3.pptx'
FIG_DIR = BASE_DIR / 'parallel-global-stiffness-assembly-research-and-implementation/parallel_global_stiffness_assembly/cpu_parallel_stiffness_assembly/results/2026-04-22/figures'
CSV_DIR = BASE_DIR / 'parallel-global-stiffness-assembly-research-and-implementation/parallel_global_stiffness_assembly/cpu_parallel_stiffness_assembly/results/2026-04-22/csv'
PLOT_SCRIPT = BASE_DIR / 'parallel-global-stiffness-assembly-research-and-implementation/parallel_global_stiffness_assembly/cpu_parallel_stiffness_assembly/scripts/plot_cpu_results.py'
JAN_PPT = BASE_DIR / '2026年01月实习生汇报/2026年01月实习生汇报-江浩华_version2.pptx'
TMP_DIR = Path('/tmp/hermes_april_v3_assets')
TMP_DIR.mkdir(parents=True, exist_ok=True)
LOGO_PATH = TMP_DIR / 'logo.png'

RED = RGBColor(156, 0, 0)
DARK = RGBColor(28, 28, 28)
GRAY = RGBColor(110, 110, 110)
LIGHT = RGBColor(245, 245, 245)
MID = RGBColor(220, 220, 220)
WHITE = RGBColor(255, 255, 255)

FONT_CN = 'Microsoft YaHei'
FONT_EN = 'Times New Roman'


def ensure_logo() -> Path:
    if LOGO_PATH.exists():
        return LOGO_PATH
    with zipfile.ZipFile(JAN_PPT) as zf:
        data = zf.read('ppt/media/image2.png')
    LOGO_PATH.write_bytes(data)
    return LOGO_PATH


def rerender_figures() -> None:
    csvs = [
        CSV_DIR / 'cube_tet4_simplified.csv',
        CSV_DIR / 'windhub_simplified.csv',
        CSV_DIR / 'windhub_physics_tet4.csv',
        CSV_DIR / 'windhub_physics_tet4_coo_sort_reduce.csv',
    ]
    cmd = ['python3', str(PLOT_SCRIPT), *map(str, csvs), '--out-dir', str(FIG_DIR)]
    subprocess.run(cmd, check=True, cwd=str(PLOT_SCRIPT.parent))


def compute_speedup_summary() -> dict[str, tuple[str, str, float, float]]:
    result = {}
    for csv_path in sorted(CSV_DIR.glob('*.csv')):
        with csv_path.open(encoding='utf-8', newline='') as f:
            rows = [r for r in csv.DictReader(f) if r.get('status') == 'PASS']
        best = max(rows, key=lambda r: float(r['speedup']))
        result[csv_path.name] = (
            best['algorithm'],
            best['threads'],
            float(best['speedup']),
            float(best.get('assembly_mean_ms') or best.get('assembly_ms') or 0.0),
        )
    return result


def add_textbox(slide, left, top, width, height, text='', size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name=FONT_CN, margin=0.05):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multi_paragraph_box(slide, left, top, width, height, paragraphs, size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()
    for i, (text, bold, color) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.name = FONT_CN
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_topbar(slide, title, page_no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_textbox(slide, Inches(0.28), Inches(0.10), Inches(5.7), Inches(0.35), title, 23, True, WHITE, PP_ALIGN.LEFT)
    logo = ensure_logo()
    slide.shapes.add_picture(str(logo), Inches(9.05), Inches(0.08), width=Inches(3.85))
    add_textbox(slide, Inches(0.28), Inches(7.02), Inches(2.0), Inches(0.18), '2026/4/24', 10, False, GRAY, font_name=FONT_EN)
    add_textbox(slide, Inches(5.70), Inches(7.02), Inches(1.9), Inches(0.18), 'Regular Report', 10, False, GRAY, PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(12.55), Inches(7.02), Inches(0.35), Inches(0.18), str(page_no), 10, False, GRAY, PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_card(slide, left, top, width, height, title, body_lines, title_bg=RED, body_bg=WHITE):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = body_bg
    outer.line.color.rgb = MID
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.52))
    header.fill.solid()
    header.fill.fore_color.rgb = title_bg
    header.line.fill.background()
    add_textbox(slide, left + Inches(0.08), top + Inches(0.10), width - Inches(0.16), Inches(0.26), title, 18, True, WHITE, PP_ALIGN.CENTER)
    paragraphs = []
    for line in body_lines:
        if isinstance(line, tuple):
            paragraphs.append(line)
        else:
            paragraphs.append((line, False, DARK))
    add_multi_paragraph_box(slide, left + Inches(0.12), top + Inches(0.64), width - Inches(0.24), height - Inches(0.76), paragraphs, 14)
    return outer


def add_caption(slide, left, top, width, text):
    add_textbox(slide, left, top, width, Inches(0.22), text, 12, False, DARK, PP_ALIGN.CENTER)


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(sldIdLst)
    slide = slides[old_index]
    sldIdLst.remove(slide)
    sldIdLst.insert(new_index, slide)


def set_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = text


def replace_text_on_slide(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text == old:
            shape.text = new
            return


def build_version3() -> None:
    rerender_figures()
    verification = compute_speedup_summary()

    copy2(VERSION2, VERSION3)
    prs = Presentation(str(VERSION3))

    # 修正 slide 5 标题，更贴近实际内容
    replace_text_on_slide(prs.slides[4], '1.3：在多核CPU平台上的并行组装策略', '1.3：实验设计与验证流程')

    # 新增 slide A: kernels
    slide_kernel = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide_kernel, '1.4：两类 kernel 的设置目的与实现逻辑', 6)
    add_textbox(slide_kernel, Inches(0.55), Inches(0.93), Inches(12.1), Inches(0.28), '为什么要同时保留 simplified kernel 与 physics_tet4 kernel：先用低成本算子验证数据通路，再用物理算子评估真实计算代价。', 18, False, DARK, PP_ALIGN.CENTER)
    add_card(
        slide_kernel,
        Inches(0.65), Inches(1.45), Inches(5.9), Inches(4.8),
        'simplified kernel',
        [
            ('定位：低成本验证算子，用来先确认 .inp 解析、CSR、scatter plan 和并行冲突处理是否正确。', False, DARK),
            ('实现逻辑：按单元自由度数直接构造一个对角占优的局部刚度矩阵；主对角项按 element_id 做微小缩放，非对角项按距离衰减。', False, DARK),
            ('特点：不依赖真实几何和材料参数，计算代价低，适合先把“数据通路”跑通。', False, DARK),
        ]
    )
    add_card(
        slide_kernel,
        Inches(6.8), Inches(1.45), Inches(5.9), Inches(4.8),
        'physics_tet4 kernel',
        [
            ('定位：更接近真实四面体有限元局部刚度计算，用来评估算法在真实算子下的并行表现。', False, DARK),
            ('实现逻辑：读取 Tet4 四个节点坐标，构造 Jacobian 与体积；求形函数导数并组装 B 矩阵，再根据材料参数构造 D 矩阵，最终计算 ke = V·BᵀDB。', False, DARK),
            ('特点：显式引入几何与材料信息，计算更重；若元素退化或求逆失败，代码会回退到 simplified kernel。', False, DARK),
        ]
    )
    add_textbox(slide_kernel, Inches(0.95), Inches(6.35), Inches(11.4), Inches(0.35), '实践顺序建议：先跑 simplified 做通路验证，再跑 physics_tet4 做真实性能评估。', 17, False, DARK, PP_ALIGN.CENTER)
    set_notes(slide_kernel, '这一页主要解释为什么我们要保留两种局部刚度算子。simplified kernel的意义在于先把输入解析、CSR结构和并行写回通路验证清楚，避免一上来就把复杂物理计算和算法性能混在一起。physics_tet4 kernel则更接近真实Tet4有限元局部刚度计算，它引入节点几何、体积、B矩阵和材料本构矩阵D，因此能够更真实地反映并行算法在实际算子下的表现。整体上，这两类kernel分别承担“通路验证”和“真实性能评估”两种角色。')

    # 新增 slide B: algorithms
    slide_algo = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide_algo, '1.5：5类并行算法简述', 7)
    add_textbox(slide_algo, Inches(0.65), Inches(0.95), Inches(11.9), Inches(0.28), '参考1月汇报中的算法介绍方式，这里把当前CPU主线中的5类并行路线做一个简短归纳。', 17, False, DARK, PP_ALIGN.CENTER)
    card_w = Inches(4.05)
    card_h = Inches(2.0)
    add_card(slide_algo, Inches(0.55), Inches(1.45), card_w, card_h, '图着色 Coloring', [
        ('先按共享节点冲突建图并做贪心着色，同色单元互不共享节点，可分批并行装配。', False, DARK),
        ('优点是避免 atomic；缺点是需要颜色预处理，颜色分布也会影响负载均衡。', False, DARK),
    ])
    add_card(slide_algo, Inches(4.64), Inches(1.45), card_w, card_h, '直接累加 Atomic', [
        ('保持自然单元遍历顺序，把局部矩阵条目直接 atomic update 到共享 CSR。', False, DARK),
        ('实现简单、额外内存最低，但高线程下会受到热点冲突和同步成本影响。', False, DARK),
    ])
    add_card(slide_algo, Inches(8.73), Inches(1.45), card_w, card_h, '线程私有 CSR', [
        ('每个线程维护一份私有 CSR values，装配时只写本地副本，最后再做归并。', False, DARK),
        ('本质上是“用内存换同步”，在真实工程网格上属于当前第一梯队。', False, DARK),
    ])
    add_card(slide_algo, Inches(2.6), Inches(3.75), card_w, card_h, '私有 COO 排序归并', [
        ('先生成 (csr_index, value) 贡献列表，再统一排序和规约。', False, DARK),
        ('研究意义强，但中间态大、排序成本高，当前更适合作为对照组。', False, DARK),
    ])
    add_card(slide_algo, Inches(6.7), Inches(3.75), card_w, card_h, '行拥有者 Row Owner', [
        ('按全局 CSR 行划分 owner，每个线程只写自己拥有的行，从而避免 atomic。', False, DARK),
        ('更接近 owner-computes 路线，但跨 owner 单元可能会重复计算局部刚度矩阵。', False, DARK),
    ])
    add_textbox(slide_algo, Inches(0.8), Inches(6.05), Inches(11.7), Inches(0.36), '说明：serial 作为唯一正确性基线和加速比基线，在这里不单列为并行路线。', 15, False, GRAY, PP_ALIGN.CENTER)
    set_notes(slide_algo, '这一页是对当前5类并行算法做一个简短说明。图着色是传统的无冲突并行思路；atomic是最直接的共享写回方案；private CSR通过线程私有副本换掉同步；COO sort-reduce代表先生成贡献再统一规约的研究型路线；row owner则更接近owner-computes。通过这一页，听众可以更清楚地理解后面结果页里这些算法名字所对应的实现逻辑。')

    # 新增 slide C: speedup gallery 1
    slide_speed1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide_speed1, '2.4：加速效果对比图总览（一）', 11)
    add_textbox(slide_speed1, Inches(0.6), Inches(0.93), Inches(12.0), Inches(0.26), '先展示三个“单 case / 单 kernel”加速效果对比图，便于直接比较线程扩展趋势。', 17, False, DARK, PP_ALIGN.CENTER)
    imgs1 = [
        ('cube_tet4_8x8x8_simplified_speedup.png', '规则网格：cube_tet4_8x8x8 + simplified'),
        ('3d-WindTurbineHub_simplified_speedup.png', '真实工程网格：WindHub + simplified'),
        ('3d-WindTurbineHub_physics_tet4_speedup.png', '真实工程网格：WindHub + physics_tet4'),
    ]
    positions = [
        (Inches(0.45), Inches(1.35), Inches(4.05), Inches(4.8)),
        (Inches(4.64), Inches(1.35), Inches(4.05), Inches(4.8)),
        (Inches(8.83), Inches(1.35), Inches(4.05), Inches(4.8)),
    ]
    for (fname, cap), (l, t, w, _h) in zip(imgs1, positions):
        slide_speed1.shapes.add_picture(str(FIG_DIR / fname), l, t, width=w)
        add_caption(slide_speed1, l, Inches(6.15), w, cap)
    set_notes(slide_speed1, '这一页集中展示三张单 case 的加速效果对比图。第一张是规则网格 cube 的 simplified 结果，第二张是真实工程网格 WindHub 在 simplified 下的结果，第三张是真实工程网格 WindHub 在 physics_tet4 下的结果。这样排列的目的，是让听众直观看到：从规则网格到真实工程网格、从低成本算子到物理算子之后，各条并行路线的加速趋势如何变化。')

    # 新增 slide D: speedup gallery 2
    slide_speed2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide_speed2, '2.5：加速效果对比图总览（二）', 12)
    add_textbox(slide_speed2, Inches(0.6), Inches(0.93), Inches(12.0), Inches(0.26), '再展示跨 case / 跨 kernel 的综合加速图。所有 speedup 图已按原始 CSV 重新生成，并与摘要结果核对。', 17, False, DARK, PP_ALIGN.CENTER)
    slide_speed2.shapes.add_picture(str(FIG_DIR / 'cross_case_best_speedup.png'), Inches(0.7), Inches(1.45), width=Inches(5.85))
    slide_speed2.shapes.add_picture(str(FIG_DIR / 'cross_kernel_best_speedup.png'), Inches(6.8), Inches(1.45), width=Inches(5.85))
    add_caption(slide_speed2, Inches(0.9), Inches(5.95), Inches(5.4), 'cross_case_best_speedup：比较不同 case 的最优加速比')
    add_caption(slide_speed2, Inches(7.0), Inches(5.95), Inches(5.4), 'cross_kernel_best_speedup：比较不同 kernel 的最优加速比')
    verify_lines = [
        f"• cube_tet4_simplified.csv：最高加速比 {verification['cube_tet4_simplified.csv'][2]:.3f}x（{verification['cube_tet4_simplified.csv'][0]} @ {verification['cube_tet4_simplified.csv'][1]} 线程）",
        f"• windhub_simplified.csv：最高加速比 {verification['windhub_simplified.csv'][2]:.3f}x（{verification['windhub_simplified.csv'][0]} @ {verification['windhub_simplified.csv'][1]} 线程）",
        f"• windhub_physics_tet4.csv：最高加速比 {verification['windhub_physics_tet4.csv'][2]:.3f}x（{verification['windhub_physics_tet4.csv'][0]} @ {verification['windhub_physics_tet4.csv'][1]} 线程）",
        f"• windhub_physics_tet4_coo_sort_reduce.csv：最高加速比 {verification['windhub_physics_tet4_coo_sort_reduce.csv'][2]:.3f}x（{verification['windhub_physics_tet4_coo_sort_reduce.csv'][0]} @ {verification['windhub_physics_tet4_coo_sort_reduce.csv'][1]} 线程）",
    ]
    box = slide_speed2.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(6.15), Inches(11.65), Inches(0.6))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = MID
    add_multi_paragraph_box(slide_speed2, Inches(1.0), Inches(6.23), Inches(11.3), Inches(0.45), [(line, False, DARK) for line in verify_lines], 11)
    set_notes(slide_speed2, '这一页展示两张综合加速效果图，分别对应跨 case 和跨 kernel 的最佳加速比比较。为了避免图和测试数据之间存在出入，我已经从原始 CSV 重新生成了所有 speedup 图，并核对了每个 CSV 的最高加速比：cube simplified 是 row_owner 在 9 线程达到 3.908x，WindHub simplified 是 row_owner 在 14 线程达到 3.695x，WindHub physics_tet4 是 private_csr 在 8 线程达到 4.686x，而 coo_sort_reduce 对照组最高只有 0.124x。这样可以确保幻灯片中引用的加速图和测试数据一致。')

    # reorder appended slides into desired positions
    # appended slides initially at end: indices 9,10,11,12 before moving (0-based)
    move_slide(prs, len(prs.slides)-4, 5)  # kernel after current slide 5
    move_slide(prs, len(prs.slides)-3, 6)  # algorithms after kernel
    move_slide(prs, len(prs.slides)-2, len(prs.slides)-1)  # speed1 before thanks placeholder position later
    move_slide(prs, len(prs.slides)-1, len(prs.slides)-1)  # no-op keeps list stable
    # After above, thanks is still last? We need move speed slides before thanks.
    # Recompute by finding thanks slide index.
    thanks_index = None
    for i, slide in enumerate(prs.slides):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
        if any('感谢您的指正' in t for t in texts):
            thanks_index = i
            break
    if thanks_index is None:
        raise RuntimeError('Cannot find thanks slide')
    # move remaining appended speed slides (currently last two slides) before thanks in order
    while len(prs.slides) - 2 >= thanks_index and any('2.4：加速效果对比图总览（一）' in getattr(shape, 'text', '') for shape in prs.slides[-2].shapes if hasattr(shape, 'text')):
        break
    # explicit locate and move by title
    def find_slide_with_title(title: str) -> int:
        for idx, s in enumerate(prs.slides):
            for sh in s.shapes:
                if hasattr(sh, 'text') and title in sh.text:
                    return idx
        raise RuntimeError(f'Cannot find slide title {title}')
    for title, insert_pos in [('2.4：加速效果对比图总览（一）', thanks_index), ('2.5：加速效果对比图总览（二）', thanks_index + 1)]:
        old = find_slide_with_title(title)
        move_slide(prs, old, insert_pos)

    prs.save(str(VERSION3))
    print(VERSION3)
    print('verification', verification)


if __name__ == '__main__':
    build_version3()
