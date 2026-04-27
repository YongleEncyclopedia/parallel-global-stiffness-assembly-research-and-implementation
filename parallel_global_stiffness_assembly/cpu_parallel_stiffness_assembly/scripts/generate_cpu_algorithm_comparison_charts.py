#!/usr/bin/env python3
"""Generate comparison charts for five CPU parallel assembly algorithms.
Data source: actual benchmark CSV under results/2026-04-22/csv/windhub_simplified.csv.
Outputs detailed and concise charts plus a one-slide summary page.
"""
from pathlib import Path
import csv, math, os, subprocess, sys
import numpy as np
import matplotlib.pyplot as plt
from matplotlib import font_manager, patheffects
from matplotlib.patches import FancyBboxPatch, Rectangle

ROOT = Path(__file__).resolve().parents[1]
CSV = ROOT / "results" / "2026-04-22" / "csv" / "windhub_simplified.csv"
OUT = ROOT / "presentation_assets" / "2026-04-cpu-algorithm-comparison"
OUT.mkdir(parents=True, exist_ok=True)

for fp in ["/System/Library/Fonts/PingFang.ttc", "/System/Library/Fonts/STHeiti Medium.ttc", "/System/Library/Fonts/Supplemental/Songti.ttc"]:
    if Path(fp).exists():
        font_manager.fontManager.addfont(fp)
        plt.rcParams["font.family"] = font_manager.FontProperties(fname=fp).get_name()
        break
plt.rcParams["axes.unicode_minus"] = False

ALGOS = [
    ("cpu_atomic", "Atomic", "#2563EB"),
    ("cpu_private_csr", "Private CSR", "#16A34A"),
    ("cpu_coo_sort_reduce", "COO Sort-Reduce", "#F97316"),
    ("cpu_graph_coloring", "Graph Coloring", "#8B5CF6"),
    ("cpu_row_owner", "Row Owner", "#E11D48"),
]
ALGO_ORDER = [a[0] for a in ALGOS]
LABEL = {a:b for a,b,c in ALGOS}
COLOR = {a:c for a,b,c in ALGOS}
BG = "#F8FAFC"; INK="#0F172A"; MUTED="#64748B"; GRID="#E2E8F0"


def load_rows():
    rows=[]
    with open(CSV, encoding="utf-8") as f:
        for r in csv.DictReader(f):
            if r["algorithm"] in ALGO_ORDER:
                rr=dict(r)
                for k in ["threads","assembly_mean_ms","total_mean_ms","speedup","efficiency","rel_l2","max_abs","extra_memory_bytes","peak_rss_mb"]:
                    rr[k] = float(rr[k])
                rr["threads"] = int(rr["threads"])
                rr["extra_memory_gb"] = rr["extra_memory_bytes"] / 1024**3
                rr["peak_rss_gb"] = rr["peak_rss_mb"] / 1024
                rows.append(rr)
    return rows

ROWS = load_rows()
THREADS = sorted({r["threads"] for r in ROWS})


def by_algo(alg):
    return sorted([r for r in ROWS if r["algorithm"] == alg], key=lambda r: r["threads"])


def best_by_speedup(alg):
    return max(by_algo(alg), key=lambda r: r["speedup"])


def setup_fig(title, subtitle, figsize=(16,9), dark=False):
    fig = plt.figure(figsize=figsize, facecolor=("#07111F" if dark else BG))
    fig.text(0.055, 0.94, title, fontsize=27, fontweight="bold", color=("white" if dark else INK), ha="left")
    fig.text(0.055, 0.902, subtitle, fontsize=13.5, color=("#C7D2FE" if dark else MUTED), ha="left")
    return fig


def style_ax(ax, title=None, dark=False):
    ax.set_facecolor("#0B1220" if dark else "white")
    if title:
        ax.set_title(title, fontsize=16, fontweight="bold", color=("white" if dark else INK), pad=12)
    ax.grid(True, color=("#25324A" if dark else GRID), lw=0.9, alpha=0.85)
    ax.set_axisbelow(True)
    for s in ax.spines.values():
        s.set_color("#334155" if dark else "#CBD5E1")
    ax.tick_params(colors=("#E5E7EB" if dark else "#334155"), labelsize=11)
    ax.xaxis.label.set_color("#E5E7EB" if dark else "#334155")
    ax.yaxis.label.set_color("#E5E7EB" if dark else "#334155")


def save(fig, name):
    png=OUT/f"{name}.png"; svg=OUT/f"{name}.svg"
    fig.savefig(png, dpi=220, facecolor=fig.get_facecolor())
    fig.savefig(svg, facecolor=fig.get_facecolor())
    plt.close(fig)
    return png, svg


def detailed_correctness():
    fig=setup_fig("Correctness Comparison / 正确性对比", "WindTurbineHub · simplified kernel · relative L2 / max absolute error vs serial CSR assembly")
    ax1=fig.add_axes([0.07,0.16,0.42,0.66]); ax2=fig.add_axes([0.56,0.16,0.37,0.66])
    for ax,t,ylabel in [(ax1,"Relative L2 error（矩阵相对范数误差）","rel_l2"),(ax2,"Max absolute error（最大绝对条目误差）","max_abs")]:
        style_ax(ax,t)
        for alg,lab,c in ALGOS:
            rs=by_algo(alg); x=[r['threads'] for r in rs]; y=[max(r[ylabel],1e-18) for r in rs]
            ax.plot(x,y,marker='o',ms=5,lw=2.2,label=lab,color=c)
        ax.set_yscale('log'); ax.set_xlabel('Threads'); ax.set_ylabel(ylabel)
        ax.set_xticks(THREADS[::2] + ([THREADS[-1]] if THREADS[-1] not in THREADS[::2] else []))
    ax1.axhspan(1e-18,1e-12,color="#DCFCE7",alpha=.55,label="near machine precision")
    ax2.axhspan(1e-18,1e-10,color="#DCFCE7",alpha=.55)
    ax1.legend(loc='upper left',fontsize=10,frameon=True,facecolor='white',edgecolor='#CBD5E1')
    fig.text(0.07,0.075,"Interpretation: all five parallel algorithms remain around floating-point round-off level; zero-error algorithms are plotted at 1e-18 for log-scale visibility.",fontsize=11.5,color=MUTED)
    return save(fig,"01_detailed_correctness")


def detailed_efficiency():
    fig=setup_fig("Efficiency Comparison / 效率对比", "Actual benchmark speedup and parallel efficiency over 1–14 CPU threads")
    ax1=fig.add_axes([0.07,0.16,0.42,0.66]); ax2=fig.add_axes([0.56,0.16,0.37,0.66])
    style_ax(ax1,"Assembly speedup vs serial baseline")
    style_ax(ax2,"Parallel efficiency = speedup / threads")
    ax1.plot(THREADS, THREADS, '--', lw=1.6, color='#94A3B8', label='Ideal linear')
    for alg,lab,c in ALGOS:
        rs=by_algo(alg); x=[r['threads'] for r in rs]
        ax1.plot(x,[r['speedup'] for r in rs],marker='o',ms=5,lw=2.4,label=lab,color=c)
        ax2.plot(x,[r['efficiency'] for r in rs],marker='o',ms=5,lw=2.4,label=lab,color=c)
    ax1.set_xlabel('Threads'); ax1.set_ylabel('Speedup ×'); ax1.set_ylim(0, max(4.1, max(r['speedup'] for r in ROWS)*1.15))
    ax2.set_xlabel('Threads'); ax2.set_ylabel('Efficiency'); ax2.set_ylim(0,1.05)
    ax1.legend(loc='upper left',fontsize=10,frameon=True,facecolor='white',edgecolor='#CBD5E1')
    fig.text(0.07,0.075,"Key observation: Row Owner reaches the highest observed speedup; Private CSR is strong but memory/reduction cost limits high-thread scaling; COO Sort-Reduce is a slow research baseline.",fontsize=11.5,color=MUTED)
    return save(fig,"02_detailed_efficiency")


def detailed_memory():
    fig=setup_fig("Memory Footprint Comparison / 内存占用对比", "Extra algorithmic memory and observed peak RSS across thread counts")
    ax1=fig.add_axes([0.07,0.16,0.42,0.66]); ax2=fig.add_axes([0.56,0.16,0.37,0.66])
    style_ax(ax1,"Extra memory requested by algorithm")
    style_ax(ax2,"Peak RSS observed during benchmark")
    for alg,lab,c in ALGOS:
        rs=by_algo(alg); x=[r['threads'] for r in rs]
        ax1.plot(x,[r['extra_memory_gb'] for r in rs],marker='o',ms=5,lw=2.4,label=lab,color=c)
        ax2.plot(x,[r['peak_rss_gb'] for r in rs],marker='o',ms=5,lw=2.4,label=lab,color=c)
    ax1.set_xlabel('Threads'); ax1.set_ylabel('Extra memory (GB)')
    ax2.set_xlabel('Threads'); ax2.set_ylabel('Peak RSS (GB)')
    ax1.legend(loc='upper left',fontsize=10,frameon=True,facecolor='white',edgecolor='#CBD5E1')
    fig.text(0.07,0.075,"Memory trade-off: Atomic and Graph Coloring are lightweight; Private CSR grows with thread count; Row Owner and COO Sort-Reduce keep large task/contribution buffers.",fontsize=11.5,color=MUTED)
    return save(fig,"03_detailed_memory")


def concise_correctness(ax):
    algs=[]; vals=[]; maxabs=[]; cols=[]
    for alg,lab,c in ALGOS:
        rs=by_algo(alg); algs.append(lab); vals.append(max(max(r['rel_l2'],1e-18) for r in rs)); maxabs.append(max(r['max_abs'] for r in rs)); cols.append(c)
    style_ax(ax,"Correctness: worst relative L2 error")
    bars=ax.bar(range(len(algs)), vals, color=cols, alpha=.9)
    ax.set_yscale('log'); ax.set_ylim(1e-18, 1e-14); ax.set_ylabel('Worst rel_L2')
    ax.set_xticks(range(len(algs)), algs, rotation=18, ha='right')
    ax.axhspan(1e-18,1e-12,color="#DCFCE7",alpha=.5)
    for b,v,ma in zip(bars,vals,maxabs):
        ax.text(b.get_x()+b.get_width()/2, v*1.4, f"{v:.1e}", ha='center', va='bottom', fontsize=9, color=INK)


def concise_efficiency(ax):
    algs=[]; vals=[]; th=[]; cols=[]
    for alg,lab,c in ALGOS:
        b=best_by_speedup(alg); algs.append(lab); vals.append(b['speedup']); th.append(b['threads']); cols.append(c)
    style_ax(ax,"Efficiency: best observed speedup")
    bars=ax.bar(range(len(algs)), vals, color=cols, alpha=.9)
    ax.set_ylabel('Speedup ×'); ax.set_ylim(0, max(vals)*1.22)
    ax.set_xticks(range(len(algs)), algs, rotation=18, ha='right')
    for b,v,t in zip(bars,vals,th): ax.text(b.get_x()+b.get_width()/2, v+0.08, f"{v:.2f}×\n@{t}T", ha='center', va='bottom', fontsize=9, color=INK)


def concise_memory(ax):
    algs=[]; vals=[]; rss=[]; cols=[]
    for alg,lab,c in ALGOS:
        b=best_by_speedup(alg); algs.append(lab); vals.append(b['extra_memory_gb']); rss.append(b['peak_rss_gb']); cols.append(c)
    style_ax(ax,"Memory: extra memory at best speedup")
    bars=ax.bar(range(len(algs)), vals, color=cols, alpha=.9)
    ax.set_ylabel('Extra memory (GB)'); ax.set_ylim(0, max(vals)*1.25+0.05)
    ax.set_xticks(range(len(algs)), algs, rotation=18, ha='right')
    for b,v in zip(bars,vals): ax.text(b.get_x()+b.get_width()/2, v+0.05, f"{v:.2f}", ha='center', va='bottom', fontsize=9, color=INK)


def simple_figures():
    outputs=[]
    for name, title, func in [
        ("04_simple_correctness", "Correctness / 正确性简洁版", concise_correctness),
        ("05_simple_efficiency", "Efficiency / 效率简洁版", concise_efficiency),
        ("06_simple_memory", "Memory / 内存简洁版", concise_memory),
    ]:
        fig=setup_fig(title, "Five CPU parallel algorithms · WindTurbineHub simplified benchmark", figsize=(16,9))
        ax=fig.add_axes([0.09,0.18,0.84,0.64])
        func(ax)
        outputs.extend(save(fig,name))
    return outputs


def summary_slide():
    fig=plt.figure(figsize=(16,9),facecolor="#07111F")
    # background gradient
    axbg=fig.add_axes([0,0,1,1]); axbg.set_axis_off(); axbg.set_xlim(0,1); axbg.set_ylim(0,1)
    X,Y=np.meshgrid(np.linspace(0,1,500),np.linspace(0,1,300))
    grad=np.zeros((300,500,3)); base=np.array([7,17,31])/255; glow=np.array([37,99,235])/255
    grad[:]=base; grad += glow*np.clip(1-((X-.15)**2+(Y-.9)**2)**0.5/.8,0,1)[...,None]*.25
    grad += np.array([225,29,72])/255*np.clip(1-((X-.9)**2+(Y-.15)**2)**0.5/.7,0,1)[...,None]*.18
    axbg.imshow(np.clip(grad,0,1),extent=[0,1,0,1],origin='lower',aspect='auto')
    fig.text(0.055,0.93,"Five CPU Parallel Algorithms: Correctness · Efficiency · Memory",fontsize=26,fontweight='bold',color='white')
    fig.text(0.055,0.895,"Actual benchmark results on 3d-WindTurbineHub / simplified kernel",fontsize=13.5,color="#CBD5E1")
    positions=[[0.055,0.16,0.28,0.62],[0.365,0.16,0.28,0.62],[0.675,0.16,0.28,0.62]]
    funcs=[concise_correctness, concise_efficiency, concise_memory]
    for pos,func in zip(positions,funcs):
        # card bg
        axc=fig.add_axes([pos[0]-0.01,pos[1]-0.015,pos[2]+0.02,pos[3]+0.04]); axc.set_axis_off()
        axc.add_patch(FancyBboxPatch((0,0),1,1,boxstyle="round,pad=0.02,rounding_size=0.04",facecolor='white',edgecolor='#94A3B8',linewidth=1.2,alpha=.96))
        ax=fig.add_axes(pos, facecolor='white')
        func(ax)
        ax.tick_params(labelsize=8.3)
        ax.title.set_fontsize(12.5)
        ax.xaxis.label.set_fontsize(9); ax.yaxis.label.set_fontsize(9)
    # takeaway strip
    axstrip=fig.add_axes([0.055,0.055,0.90,0.065]); axstrip.set_axis_off()
    axstrip.add_patch(FancyBboxPatch((0,0),1,1,boxstyle="round,pad=0.02,rounding_size=0.04",facecolor="#0F172A",edgecolor="#334155",linewidth=1.2,alpha=.92))
    axstrip.text(0.02,0.53,"Takeaway",fontsize=13.5,color="#93C5FD",fontweight='bold',va='center')
    axstrip.text(0.15,0.53,"All algorithms pass correctness; Row Owner is fastest; Atomic/Graph Coloring are memory-light, while Private CSR and COO/Owner trade memory for conflict avoidance.",fontsize=12.3,color='white',va='center')
    return save(fig,"07_summary_slide_three_simple_charts")


def write_data_summary():
    lines=[]
    lines.append("# CPU parallel algorithm comparison data summary\n")
    lines.append(f"Data source: `{CSV}`\n")
    lines.append("Case: 3d-WindTurbineHub, kernel: simplified. Five parallel algorithms are compared against cpu_serial.\n")
    lines.append("| Algorithm | Best speedup | Threads | Assembly ms | Total ms | Worst rel_L2 | Worst max_abs | Extra memory at best (GB) | Peak RSS at best (GB) |\n")
    lines.append("|---|---:|---:|---:|---:|---:|---:|---:|---:|\n")
    for alg,lab,c in ALGOS:
        b=best_by_speedup(alg); rs=by_algo(alg)
        lines.append(f"| {lab} | {b['speedup']:.3f} | {b['threads']} | {b['assembly_mean_ms']:.2f} | {b['total_mean_ms']:.2f} | {max(r['rel_l2'] for r in rs):.3e} | {max(r['max_abs'] for r in rs):.3e} | {b['extra_memory_gb']:.3f} | {b['peak_rss_gb']:.2f} |\n")
    p=OUT/"data_summary.md"; p.write_text(''.join(lines),encoding='utf-8')
    return p


def make_pptx(summary_png):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        subprocess.check_call([sys.executable,"-m","pip","install","python-pptx","-q"])
        from pptx import Presentation
        from pptx.util import Inches
    prs=Presentation(); prs.slide_width=Inches(16); prs.slide_height=Inches(9)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(summary_png),0,0,width=prs.slide_width,height=prs.slide_height)
    pptx=OUT/"cpu_algorithm_comparison_summary_slide.pptx"
    prs.save(pptx)
    return pptx

if __name__ == "__main__":
    outs=[]
    outs += list(detailed_correctness())
    outs += list(detailed_efficiency())
    outs += list(detailed_memory())
    outs += simple_figures()
    summary_png, summary_svg = summary_slide(); outs += [summary_png, summary_svg]
    outs.append(write_data_summary())
    outs.append(make_pptx(summary_png))
    print("Generated comparison assets:")
    for p in outs: print(p)
