from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pathlib import Path

OUT_DIR = Path('/Users/macstudio/Documents/Intern_Peking University_supu/2026年04月实习生汇报')
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPT = OUT_DIR / '2026年04月实习生汇报-江浩华.pptx'
OUT_NOTES = OUT_DIR / '2026年04月实习生汇报-江浩华_逐页讲解文案.md'

ASSET_DIR = Path('/tmp/jhh_ppt_assets')
BANNER = ASSET_DIR / 'image4.jpeg'
LOGO = ASSET_DIR / 'image2.png'
FIG_DIR = Path('/Users/macstudio/Documents/Intern_Peking University_supu/parallel-global-stiffness-assembly-research-and-implementation/parallel_global_stiffness_assembly/cpu_parallel_stiffness_assembly/results/2026-04-22/figures')
SIMPLIFIED_DASH = FIG_DIR / '3d-WindTurbineHub_simplified_dashboard.png'
PHYSICS_DASH = FIG_DIR / '3d-WindTurbineHub_physics_tet4_dashboard.png'
CROSS_KERNEL = FIG_DIR / 'cross_kernel_best_speedup.png'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

RED = RGBColor(154, 0, 0)
DARK = RGBColor(25, 25, 25)
GRAY = RGBColor(120, 120, 120)
LIGHT = RGBColor(245, 245, 245)
MID = RGBColor(230, 230, 230)
WHITE = RGBColor(255, 255, 255)

FONT_CN = 'PingFang SC'
FONT_EN = 'Times New Roman'


def add_textbox(slide, left, top, width, height, text='', font_size=20, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name=FONT_CN, line_spacing=1.2):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    p.line_spacing = line_spacing
    return tx


def add_bullet_lines(slide, left, top, width, lines, font_size=20, level0_color=DARK, bullet_color=RED):
    tx = slide.shapes.add_textbox(left, top, width, Inches(3.8))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = level0_color
        run.font.name = FONT_CN
        if line.startswith('•'):
            run.font.color.rgb = DARK
        elif line.startswith('➢'):
            run.font.color.rgb = RED
            run.font.bold = True
    return tx


def add_topbar(slide, title, page_no):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_textbox(slide, Inches(0.35), Inches(0.12), Inches(5.5), Inches(0.4), f'▪ {title}', 24, True, WHITE)
    slide.shapes.add_picture(str(LOGO), Inches(9.1), Inches(0.09), height=Inches(0.5))
    add_textbox(slide, Inches(0.25), Inches(7.05), Inches(2.5), Inches(0.22), '2026/4/30', 10, False, GRAY, font_name=FONT_EN)
    add_textbox(slide, Inches(5.75), Inches(7.05), Inches(1.8), Inches(0.22), 'Regular Report', 10, False, GRAY, PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(12.6), Inches(7.02), Inches(0.35), Inches(0.22), str(page_no), 11, False, GRAY, PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(BANNER), 0, 0, width=prs.slide_width, height=Inches(2.95))
    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.2), Inches(2.18), Inches(13.8), Inches(1.65))
    arc.fill.solid()
    arc.fill.fore_color.rgb = WHITE
    arc.line.fill.background()
    slide.shapes.add_picture(str(LOGO), Inches(3.2), Inches(2.35), width=Inches(6.9))
    add_textbox(slide, Inches(1.2), Inches(4.05), Inches(11), Inches(0.45), '2026年04月工作汇报：', 23, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.85), Inches(4.55), Inches(11.7), Inches(0.9), 'CPU并行整体刚度矩阵组装\n可复现实验平台建设与阶段验证', 26, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.25), Inches(6.05), Inches(4.8), Inches(0.55), '报告人：江浩华    合作导师：苏璞', 17, False, DARK, PP_ALIGN.CENTER)
    return slide


def add_section_slide(title, subtitle, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_picture(str(LOGO), Inches(9.1), Inches(0.45), width=Inches(3.7))
    add_textbox(slide, Inches(0.95), Inches(4.18), Inches(6.5), Inches(0.7), title, 28, True, DARK)
    add_textbox(slide, Inches(0.95), Inches(4.88), Inches(5.5), Inches(0.45), subtitle, 18, False, GRAY)
    add_textbox(slide, Inches(0.25), Inches(7.05), Inches(2.5), Inches(0.22), '2026/4/30', 10, False, GRAY, font_name=FONT_EN)
    add_textbox(slide, Inches(5.75), Inches(7.05), Inches(1.8), Inches(0.22), 'Regular Report', 10, False, GRAY, PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(12.6), Inches(7.02), Inches(0.35), Inches(0.22), str(page_no), 11, False, GRAY, PP_ALIGN.RIGHT, font_name=FONT_EN)
    return slide


def add_background_goals(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, '1.1：项目背景与本月目标', page_no)
    add_textbox(slide, Inches(0.6), Inches(1.0), Inches(6.0), Inches(0.35), '项目背景', 22, True, RED)
    lines = [
        '➢ 1月阶段：完成多类并行组装策略调研，整体叙事仍偏 GPU-first。',
        '➢ 4月阶段：项目目标转向 CPU 可复现实验平台建设。',
        '• 在统一框架下比较正确性、性能、扩展性与内存代价。',
        '• 在规则网格与真实工程网格上输出统一 benchmark 结果。',
    ]
    add_bullet_lines(slide, Inches(0.7), Inches(1.42), Inches(5.8), lines, 17)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.18), Inches(5.6), Inches(4.9))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = MID
    add_textbox(slide, Inches(7.25), Inches(1.45), Inches(5.0), Inches(0.35), '本月目标', 21, True, RED)
    goal_lines = [
        '1. 固定唯一 CPU 主线目录，统一核心文档。',
        '2. 跑通真实工程网格 3d-WindTurbineHub 实验矩阵。',
        '3. 输出 CSV / JSON / Markdown / PNG / SVG 五类结果。',
        '4. 基于真实数据形成阶段性算法判断。',
    ]
    add_bullet_lines(slide, Inches(7.25), Inches(1.9), Inches(4.95), [f'• {x}' for x in goal_lines], 16)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(5.25), Inches(5.1), Inches(0.55))
    band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background()
    add_textbox(slide, Inches(7.35), Inches(5.36), Inches(4.8), Inches(0.25), '项目定义：从算法原型推进为可复现实验平台', 15, True, WHITE, PP_ALIGN.CENTER)


def add_platform_build(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, '1.2：平台建设与本月完成内容', page_no)
    titles = ['主线收敛', '实验框架', '结果产物']
    bullets = [
        ['唯一主线目录固定为 cpu_parallel_stiffness_assembly', '历史 GPU 资产隔离为 legacy 参考', 'README / requirements / handoff 口径一致'],
        ['统一算法工厂与 benchmark CLI', '统一规则网格、.inp 解析、CSR 与 scatter plan', '支持 serial / atomic / private_csr / coo_sort_reduce / coloring / row_owner'],
        ['自动输出 CSV / JSON / Markdown / PNG / SVG', '图表直接面向论文/PPT 使用', '沉淀可复现脚本与验证记录'],
    ]
    xs = [0.55, 4.52, 8.49]
    for i in range(3):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs[i]), Inches(1.3), Inches(3.6), Inches(4.5))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = MID
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xs[i]), Inches(1.3), Inches(3.6), Inches(0.62))
        head.fill.solid(); head.fill.fore_color.rgb = RED; head.line.fill.background()
        add_textbox(slide, Inches(xs[i]+0.12), Inches(1.47), Inches(3.25), Inches(0.2), titles[i], 19, True, WHITE, PP_ALIGN.CENTER)
        add_bullet_lines(slide, Inches(xs[i]+0.18), Inches(2.05), Inches(3.05), [f'• {x}' for x in bullets[i]], 16)
    add_textbox(slide, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.5), '本月工作的核心不是“新增零散功能”，而是完成主线统一、实验补齐、结果固化与结论提炼四项平台化建设。', 17, False, DARK, PP_ALIGN.CENTER)


def add_experiment_design(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, '1.3：实验设计与验证流程', page_no)
    add_textbox(slide, Inches(0.6), Inches(0.98), Inches(5.8), Inches(0.32), '统一实验流程', 22, True, RED)
    steps = ['统一输入\n规则网格 + Abaqus .inp', '统一算法\n6类 CPU 装配路线', '统一运行\n脚本化 benchmark', '统一产物\nCSV/JSON/图表/摘要']
    x = 0.7
    for idx, s in enumerate(steps):
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(2.6), Inches(0.9))
        r.fill.solid(); r.fill.fore_color.rgb = LIGHT; r.line.color.rgb = MID
        add_textbox(slide, Inches(x+0.1), Inches(1.85), Inches(2.4), Inches(0.45), s, 17, True, DARK, PP_ALIGN.CENTER)
        if idx < len(steps)-1:
            add_textbox(slide, Inches(x+2.63), Inches(1.93), Inches(0.35), Inches(0.2), '→', 24, True, RED, PP_ALIGN.CENTER)
        x += 2.95
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.1), Inches(5.4), Inches(2.9))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.color.rgb = MID
    add_textbox(slide, Inches(0.95), Inches(3.35), Inches(4.9), Inches(0.28), '真实工程网格规模', 21, True, RED)
    stat_lines = [
        '• 算例：3d-WindTurbineHub.inp',
        '• 节点数：228,384',
        '• 单元数：1,113,684',
        '• 总自由度：685,152',
        '• CSR NNZ：27,502,200',
    ]
    add_bullet_lines(slide, Inches(1.02), Inches(3.78), Inches(4.6), stat_lines, 17)
    slide.shapes.add_picture(str(CROSS_KERNEL), Inches(6.5), Inches(1.35), width=Inches(5.9))
    add_textbox(slide, Inches(6.7), Inches(5.65), Inches(5.5), Inches(0.42), '统一输出 cross-kernel / cross-case 图表，为汇报与论文直接复用提供基础。', 16, False, DARK, PP_ALIGN.CENTER)


def add_result_slide(page_no, title, image_path, bullets, label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, title, page_no)
    add_textbox(slide, Inches(0.62), Inches(0.98), Inches(4.6), Inches(0.3), label, 21, True, RED)
    add_bullet_lines(slide, Inches(0.75), Inches(1.45), Inches(4.2), bullets, 16)
    slide.shapes.add_picture(str(image_path), Inches(4.95), Inches(1.05), width=Inches(7.9))


def add_comparison_slide(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, '2.3：算法对比与阶段性判断', page_no)
    add_textbox(slide, Inches(0.6), Inches(0.98), Inches(5.2), Inches(0.32), '当前算法优先级', 22, True, RED)
    headers = ['算法路线', '速度表现', '资源代价', '当前判断']
    rows = [
        ['row_owner', 'simplified 最优\n55.141 ms / 3.695x', '1.792 GiB', '高线程表现最强\n主推路线'],
        ['private_csr', 'physics_tet4 最优\n119.566 ms / 4.686x', '1.639 GiB', '中高线程稳定\n主推路线'],
        ['atomic', '接近第一梯队\n123.174 ms / 4.549x', '0 GiB', '低内存代价\n工程性价比高'],
        ['coloring', '稳定可运行\n但已不领先', '0.008 GiB', '保留为重要基线'],
        ['coo_sort_reduce', '正确但显著偏慢\n0.109x~0.124x', '2.390 GiB', '研究对照组\n不作主推'],
    ]
    left = Inches(0.62); top = Inches(1.42)
    widths = [Inches(2.05), Inches(2.75), Inches(1.55), Inches(3.9)]
    row_h = Inches(0.70)
    x = left
    for i, h in enumerate(headers):
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, widths[i], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = RED; cell.line.color.rgb = WHITE
        add_textbox(slide, x+Inches(0.03), top+Inches(0.18), widths[i]-Inches(0.06), Inches(0.22), h, 16, True, WHITE, PP_ALIGN.CENTER)
        x += widths[i]
    for r, row in enumerate(rows):
        x = left
        y = top + row_h*(r+1)
        for c, text in enumerate(row):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, widths[c], row_h)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT
            cell.line.color.rgb = MID
            add_textbox(slide, x+Inches(0.05), y+Inches(0.12), widths[c]-Inches(0.1), Inches(0.45), text, 13.5, False, DARK, PP_ALIGN.CENTER)
            x += widths[c]
    add_textbox(slide, Inches(0.85), Inches(6.08), Inches(11.55), Inches(0.42), '阶段结论：当前最值得继续深入的是 row_owner 与 private_csr；atomic 作为低内存代价方案，也具备很强工程价值。', 16, False, DARK, PP_ALIGN.CENTER)


def add_summary_slide(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_topbar(slide, '3.1：阶段总结与下一步计划', page_no)
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(5.7), Inches(4.9))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = WHITE; left_box.line.color.rgb = MID
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.2), Inches(5.7), Inches(4.9))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = LIGHT; right_box.line.color.rgb = MID
    add_textbox(slide, Inches(0.95), Inches(1.45), Inches(5.1), Inches(0.28), '本月阶段性成果', 21, True, RED)
    add_bullet_lines(slide, Inches(1.0), Inches(1.92), Inches(4.9), [
        '• 完成 CPU-first 叙事切换，明确唯一主线。',
        '• 建成可复现实验平台，沉淀脚本、文档与统一输出。',
        '• 在真实工程网格上跑通多算法统一 benchmark。',
        '• 形成“主推路线 + 工程折中方案 + 对照组”的阶段判断。',
    ], 16)
    add_textbox(slide, Inches(7.15), Inches(1.45), Inches(5.1), Inches(0.28), '下一阶段计划', 21, True, RED)
    add_bullet_lines(slide, Inches(7.2), Inches(1.92), Inches(4.95), [
        '• 扩展 physics_tet4 在线程维度上的完整扫描。',
        '• 在 Windows + Intel 平台复现实验矩阵。',
        '• 继续优化图表版式与论文/PPT 级摘要输出。',
        '• 围绕 row_owner / private_csr / atomic 进一步做性能归因分析。',
    ], 16)
    add_textbox(slide, Inches(0.9), Inches(6.25), Inches(11.7), Inches(0.36), '本月工作的意义：不是单独实现某一个算法，而是把 CPU 并行整体刚度矩阵研究推进成了可验证、可比较、可复现的工程级研究平台。', 17, False, DARK, PP_ALIGN.CENTER)


def add_thanks(page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_picture(str(LOGO), Inches(9.15), Inches(0.5), width=Inches(3.55))
    add_textbox(slide, Inches(3.3), Inches(3.1), Inches(6.5), Inches(0.7), '感谢您的指正', 30, True, DARK, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.25), Inches(7.05), Inches(2.5), Inches(0.22), '2026/4/30', 10, False, GRAY, font_name=FONT_EN)
    add_textbox(slide, Inches(5.75), Inches(7.05), Inches(1.8), Inches(0.22), 'Regular Report', 10, False, GRAY, PP_ALIGN.CENTER, font_name=FONT_EN)
    add_textbox(slide, Inches(12.6), Inches(7.02), Inches(0.35), Inches(0.22), str(page_no), 11, False, GRAY, PP_ALIGN.RIGHT, font_name=FONT_EN)


add_title_slide()
add_section_slide('CPU并行整体刚度矩阵', '可复现实验平台建设与阶段验证', 2)
add_background_goals(3)
add_platform_build(4)
add_experiment_design(5)
add_result_slide(6, '2.1：真实工程网格结果（一）', SIMPLIFIED_DASH, [
    '➢ 工程网格 + simplified kernel：row_owner 在 14 线程达到当前最优。',
    '• 最优结果：55.141 ms，3.695x 加速。',
    '• private_csr 次优：71.211 ms，2.861x；atomic 以 0 GiB 额外内存达到 80.560 ms。',
    '• 说明：在 simplified 场景下，owner-computes 路线的高线程优势最明显。',
], '真实工程网格结果：WindHub + simplified')
add_result_slide(7, '2.2：真实工程网格结果（二）', PHYSICS_DASH, [
    '➢ 工程网格 + physics_tet4 kernel：private_csr 与 row_owner 构成第一梯队。',
    '• 最优结果：private_csr @ 8 线程，119.566 ms，4.686x。',
    '• row_owner @ 8 线程几乎持平：120.444 ms，4.652x；atomic 也达到 123.174 ms，4.549x。',
    '• 说明：更高物理复杂度下，私有缓冲归并路线的稳定性优势更加明显。',
], '真实工程网格结果：WindHub + physics_tet4')
add_comparison_slide(8)
add_summary_slide(9)
add_thanks(10)

prs.save(OUT_PPT)

notes = '''# 2026年04月实习生汇报——逐页讲解文案

## 第1页 封面页
本次汇报的主题是“CPU并行整体刚度矩阵组装可复现实验平台建设与阶段验证”。
和1月阶段相比，本月工作的重点已经从并行组装算法的前期调研，进一步推进到面向多核CPU的实验平台建设、真实工程网格验证以及阶段性算法判断。

## 第2页 章节过渡页
这一部分主要汇报两方面内容：第一，为什么本月工作要从“原型探索”转向“平台建设”；第二，在真实工程网格上，我们已经形成了哪些可复现、可比较的阶段成果。

## 第3页 项目背景与本月目标
1月阶段已经完成了多类并行组装策略调研，但当时的叙事仍然偏向GPU-first。
4月阶段的核心变化，是把项目重新定义为CPU并行整体刚度矩阵组装的可复现实验平台建设。
本月的四个目标分别是：固定唯一CPU主线目录、补齐真实工程网格实验矩阵、输出结构化实验结果，以及基于统一数据形成阶段结论。

## 第4页 平台建设与本月完成内容
本月完成内容可以概括为三类。
第一类是主线收敛：明确唯一主线目录，历史GPU资产转为legacy参考，并统一README、需求文档和交接文档口径。
第二类是实验框架：统一算法工厂、benchmark CLI、输入解析和CSR/scatter plan流程，使不同算法能够在同一口径下比较。
第三类是结果产物：实验现在能够自动输出CSV、JSON、Markdown、PNG、SVG五类结果，图表可以直接服务于PPT和论文材料。

## 第5页 实验设计与验证流程
本月实验设计强调“统一流程”。
即统一输入、统一算法、统一运行方式、统一输出结果。
在实验对象上，我们既保留了规则网格用于回归，也重点引入了真实工程网格3d-WindTurbineHub。
这个工程算例的规模已经达到22.8万个节点、111.4万个单元、68.5万个自由度和2750万级别NNZ，因此它足以支撑工程级比较，而不是停留在玩具样例层面。

## 第6页 真实工程网格结果（一）——simplified
在真实工程网格、simplified kernel下，当前最优算法是row_owner。
它在14线程下达到了55.141毫秒，对串行基线实现了3.695倍加速。
与此同时，private_csr达到71.211毫秒、2.861倍加速；atomic在额外内存几乎为零的前提下达到80.560毫秒。
这说明在simplified场景下，owner-computes路线在高线程条件下优势最明显。

## 第7页 真实工程网格结果（二）——physics_tet4
在更接近真实计算代价的physics_tet4 kernel下，第一梯队收敛为private_csr与row_owner。
其中最优结果是private_csr在8线程达到119.566毫秒、4.686倍加速。
row_owner在8线程的结果几乎持平，为120.444毫秒、4.652倍；atomic也达到123.174毫秒、4.549倍。
这表明在更复杂kernel下，线程私有缓冲归并路线的稳定性更突出，而row_owner仍保持非常强的竞争力。

## 第8页 算法对比与阶段性判断
基于统一平台下的实验结果，当前最值得继续深入的两类路线已经比较清楚：row_owner和private_csr。
其中，row_owner在simplified场景下表现最强，private_csr在physics_tet4场景下略优。
atomic虽然不是绝对最快，但由于额外内存几乎为零，具有非常高的工程性价比。
coloring可以稳定运行，但已经不是当前领先路线；coo_sort_reduce虽然正确性成立，但速度和资源代价都不具备优势，更适合作为研究对照组。

## 第9页 阶段总结与下一步计划
本月工作的意义，不是单独实现了某一个并行算法，而是把CPU并行整体刚度矩阵研究推进成了一套可验证、可比较、可复现的工程级实验平台。
下一阶段主要有四项工作：扩展physics_tet4的完整线程扫描；在Windows+Intel平台复现实验矩阵；继续优化图表与论文/PPT级摘要输出；围绕row_owner、private_csr和atomic进一步做性能归因分析。

## 第10页 结束页
以上是本月汇报内容，感谢老师和各位的指正。
'''
OUT_NOTES.write_text(notes, encoding='utf-8')
print(OUT_PPT)
print(OUT_NOTES)
